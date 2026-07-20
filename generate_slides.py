import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Use 16:9 widescreen layout (standard modern presentation size)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom Theme Palette (Emerald Dark Mode matching the CropDiag UI)
    bg_color = RGBColor(12, 15, 29)       # #0C0F1D (Dark navy/slate)
    text_white = RGBColor(243, 244, 246)  # #F3F4F6 (Off-white)
    text_gray = RGBColor(156, 163, 175)   # #9CA3AF (Slate gray)
    emerald_green = RGBColor(16, 185, 129) # #10B981 (Bright Emerald)
    red_accent = RGBColor(239, 68, 68)    # #EF4444 (Crimson)

    # Helper function to apply background to slides
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Helper function to add headers
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = emerald_green
        
        # Border accent line
        line_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.833), Inches(0.05))
        tf_line = line_box.text_frame
        p_line = tf_line.paragraphs[0]
        p_line.text = "_________________________________________________________________________________________________"
        p_line.font.name = "Arial"
        p_line.font.size = Pt(10)
        p_line.font.bold = True
        p_line.font.color.rgb = emerald_green

    # ------------------ SLIDE 1: COVER PAGE ------------------
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1)
    
    # College Title
    clg_box = slide1.shapes.add_textbox(Inches(1), Inches(0.6), Inches(11.333), Inches(1.2))
    tf = clg_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "V.S.B ENGINEERING COLLEGE (AUTONOMOUS), KARUR - 639111"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = emerald_green
    
    p2 = tf.add_paragraph()
    p2.text = "DEPARTMENT OF INFORMATION TECHNOLOGY"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Arial"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = text_white

    # Project Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AUTOMATED CROP DISEASE ANALYSIS USING CNN ALGORITHM"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_white
    
    p_sub = tf.add_paragraph()
    p_sub.text = "A Full-Stack Deep Learning Application for Agricultural Pathology Detection"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = text_gray

    # Candidate info
    info_box = slide1.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.333), Inches(1.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Submitted by:"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.color.rgb = text_gray
    
    p2 = tf.add_paragraph()
    p2.text = "MOUNISHA P"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = emerald_green
    
    p3 = tf.add_paragraph()
    p3.text = "Register Number: 922524243113"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = text_white

    # Footer date
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ACADEMIC YEAR: 2025 - 2026 | JULY 2026"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = text_gray

    # ------------------ SLIDE 2: INTRODUCTION ------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2)
    add_slide_header(slide2, "1. Project Introduction")
    
    content_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Overview of the Problem:"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = text_white
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "• Crop leaf diseases annually degrade up to 30% of global yield capacities, threatening agrarian livelihoods.\n" \
              "• Farmers often lack access to professional agronomists, leading to incorrect diagnostic treatments.\n" \
              "• Traditional identification methods are manual, slow, subjective, and highly error-prone."
    p2.font.name = "Arial"
    p2.font.size = Pt(16)
    p2.font.color.rgb = text_gray
    p2.space_after = Pt(24)
    
    p3 = tf.add_paragraph()
    p3.text = "Deep Learning Solution:"
    p3.font.name = "Arial"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = text_white
    p3.space_after = Pt(10)
    
    p4 = tf.add_paragraph()
    p4.text = "• Integrates Computer Vision (OpenCV HSV Masks) with Deep Learning architectures (CNN / MLP).\n" \
              "• Provides instant automated diagnostics for 10 distinct crops, split into 28 healthy and diseased categories.\n" \
              "• Issues real-time treatment recipes (organic, chemical, water, climate) to mitigate crop damage."
    p4.font.name = "Arial"
    p4.font.size = Pt(16)
    p4.font.color.rgb = text_gray

    # ------------------ SLIDE 3: OBJECTIVES ------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3)
    add_slide_header(slide3, "2. Project Objectives")
    
    content_box = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    objectives = [
        ("Early Detection", "Diagnose pathogens at initial leaf-spot stages using color space analysis to limit outbreak spreads."),
        ("High-Precision Accuracy", "Leverage deep neural layers to achieve high validation scores on standard agricultural datasets."),
        ("Accessible Actionable Advice", "Provide comprehensive treatment summaries, specifying organic recipes and chemical control values."),
        ("Robust Portability & Execution", "Build a standalone executable that works offline without needing external servers, complex configurations, or terminal windows."),
        ("Academic Compliance", "Provide structured PDF reports with dynamic verification QR codes, matching final-year review requirements.")
    ]
    
    for title, desc in objectives:
        p = tf.add_paragraph()
        p.text = f"✔ {title}: "
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = emerald_green
        
        # Append description in white
        p_desc = p.add_run()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(15)
        p_desc.font.bold = False
        p_desc.font.color.rgb = text_white
        p.space_after = Pt(14)

    # ------------------ SLIDE 4: TECHNOLOGY STACK ------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4)
    add_slide_header(slide4, "3. Technology Stack & Frameworks")
    
    # 4 Columns represented as textboxes
    cols = [
        ("Frontend", "React.js & TypeScript", "• Component layout isolation\n• Compile-time type safety\n• Single page routing\n• Lucid-React Icons"),
        ("Styling", "Vanilla CSS & Glassmorphism", "• Custom emerald themes\n• Tailored dark palettes\n• Hover micro-animations\n• Responsive layouts"),
        ("Backend Gateway", "Python Flask", "• High-speed HTTP services\n• Multi-part image parser\n• Static asset dispatcher\n• Dynamic PDF compiler"),
        ("Database & ML", "SQLite & Scikit-Learn", "• Portability without db setups\n• Parametric SQL security\n• MLP classifier model\n• OpenCV HSV filter masks")
    ]
    
    for i, (name, frame, features) in enumerate(cols):
        col_box = slide4.shapes.add_textbox(Inches(0.75 + i*2.95), Inches(2.0), Inches(2.8), Inches(4.5))
        tf = col_box.text_frame
        tf.word_wrap = True
        
        # Title
        p_name = tf.paragraphs[0]
        p_name.text = name
        p_name.font.name = "Arial"
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = emerald_green
        p_name.space_after = Pt(6)
        
        # Framework
        p_frame = tf.add_paragraph()
        p_frame.text = frame
        p_frame.font.name = "Arial"
        p_frame.font.size = Pt(14)
        p_frame.font.bold = True
        p_frame.font.color.rgb = text_white
        p_frame.space_after = Pt(14)
        
        # Features
        p_features = tf.add_paragraph()
        p_features.text = features
        p_features.font.name = "Arial"
        p_features.font.size = Pt(12)
        p_features.font.color.rgb = text_gray

    # ------------------ SLIDE 5: SYSTEM ARCHITECTURE ------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5)
    add_slide_header(slide5, "4. System Architecture")
    
    # Left side: Text explanation
    left_box = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(5.0))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Core Architectural Flow:"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(12)
    
    arch_points = [
        ("Client Tier", "React SPA collects leaf image uploads, communicates with the API Gateway using Axios async promises, and downloads Compiled PDF reports."),
        ("Application Tier", "Flask gateway receives images, forwards raw bytes into OpenCV preprocessing masks, extracts features, and runs diagnostic predictions through the MLP Model."),
        ("Data Tier", "SQLite commits records into normalization schemas. ReportLab builds a university-grade PDF with vector QR codes and returns it dynamically.")
    ]
    
    for tier, desc in arch_points:
        p_tier = tf.add_paragraph()
        p_tier.text = f"• {tier}: "
        p_tier.font.name = "Arial"
        p_tier.font.size = Pt(14)
        p_tier.font.bold = True
        p_tier.font.color.rgb = emerald_green
        
        p_desc = p_tier.add_run()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = text_gray
        p_tier.space_after = Pt(12)

    # Right side: Visual Diagram representation using simple text boxes
    right_box = slide5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.0))
    tf_diag = right_box.text_frame
    tf_diag.word_wrap = True
    
    diagram_str = "========================================\n" \
                  "       CLIENT WEB PRESENTATION TIER\n" \
                  "          React.js / TypeScript SPA\n" \
                  "========================================\n" \
                  "                   ⬇️ (HTTP / Multipart Form Data)\n" \
                  "========================================\n" \
                  "        APPLICATION BUSINESS LOGIC TIER\n" \
                  "             Python Flask API Gateway\n" \
                  "   OpenCV Masking ➡️ Neural Network Inference\n" \
                  "========================================\n" \
                  "         ⬇️                          ⬇️\n" \
                  "===================         ====================\n" \
                  "  DATA TIER (SQLite)         REPORTLAB SERVICE\n" \
                  "   Parametric Queries         Structured PDFs\n" \
                  "===================         ===================="
                  
    p_diag = tf_diag.paragraphs[0]
    p_diag.text = diagram_str
    p_diag.font.name = "Consolas"
    p_diag.font.size = Pt(13)
    p_diag.font.color.rgb = emerald_green

    # ------------------ SLIDE 6: MODULES ------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6)
    add_slide_header(slide6, "5. Project Modules")
    
    content_box = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    modules = [
        ("Module 1: User Session Security", "Manages credentials, SHA-256 hashed password registries, and token-based API authentication controls."),
        ("Module 2: OpenCV Preprocessing", "Applies HSV Color space filtering to isolates healthy leaf green regions from yellow/brown spots, computing precise damage percentages."),
        ("Module 3: Neural Network Diagnosis", "Inferences the Scikit-Learn Multilayer Perceptron (MLP) Classifier to diagnose 28 specific agricultural healthy and diseased target classes."),
        ("Module 4: SQL Database Registry", "Records histories of prediction indexes, confidence ratios, date, and user IDs in normalized tables."),
        ("Module 5: Dynamic Report Generation", "Utilizes ReportLab flowables to build verified PDF reports with vector QR codes pointing back to host diagnostics.")
    ]
    
    for title, desc in modules:
        p = tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = emerald_green
        
        p_desc = p.add_run()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = text_white
        p.space_after = Pt(12)

    # ------------------ SLIDE 7: PREPROCESSING & OPENCV ------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7)
    add_slide_header(slide7, "6. Image Preprocessing & HSV Masking")
    
    content_box = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "How OpenCV Enhances CNN Inference:"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(12)
    
    cv_steps = [
        ("Color Space Conversion", "Converts raw RGB input pixels to Hue-Saturation-Value (HSV). The HSV space is highly resistant to outdoor lighting variances, shadow castings, and glare reflections."),
        ("Green Mask Filtering", "Defines HSV threshold limits for green chlorophyllic structures (Hue: 35-85, Saturation: 40-255). It isolates the true leaf blade contours from pot backgrounds or soil."),
        ("Necrotic Spot Masking", "Defines thresholds for brown, black, and yellow dead spots. Counts non-green pixels inside the leaf border to compute the exact damage ratio."),
        ("Pathology Classification Override", "If the leaf features less than 1% total necrotic spots, it overrides neural predictions to return a clean 'Healthy' status, eliminating false positive detections.")
    ]
    
    for title, desc in cv_steps:
        p_step = tf.add_paragraph()
        p_step.text = f"✔ {title}: "
        p_step.font.name = "Arial"
        p_step.font.size = Pt(15)
        p_step.font.bold = True
        p_step.font.color.rgb = emerald_green
        
        p_desc = p_step.add_run()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = text_gray
        p_step.space_after = Pt(10)

    # ------------------ SLIDE 8: CNN MODEL & FALLBACK ------------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide8)
    add_slide_header(slide8, "7. Deep Neural Network Classifier")
    
    content_box = slide8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Neural Network Pipeline Specs:"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(12)
    
    nn_specs = [
        ("Input Feature Vectorization", "Leaf photos are loaded and downsampled to 32x32 pixel maps, vectorized into 3072 raw input features (RGB channels) for fast CPU inference speeds."),
        ("Multi-Layer Perceptron (MLP)", "Configures input layers, hidden layer 1 (128 neurons), hidden layer 2 (64 neurons), and an output Softmax layer mapping to 28 target diagnostic classes."),
        ("Activation & Optimizations", "Uses Rectified Linear Units (ReLU) in hidden layers to capture complex nonlinear spot patterns, optimizing loss through Stochastic Gradient Descent (SGD)."),
        ("CNN Training Setup", "The system natively runs Scikit-Learn MLPClassifier when TensorFlow dependencies lack AVX CPU instruction sets, loading real PlantVillage photos to save weights into disease_model.h5.")
    ]
    
    for title, desc in nn_specs:
        p_nn = tf.add_paragraph()
        p_nn.text = f"• {title}: "
        p_nn.font.name = "Arial"
        p_nn.font.size = Pt(15)
        p_nn.font.bold = True
        p_nn.font.color.rgb = emerald_green
        
        p_desc = p_nn.add_run()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = text_gray
        p_nn.space_after = Pt(10)

    # ------------------ SLIDE 9: SQL DATABASE DESIGN ------------------
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide9)
    add_slide_header(slide9, "8. SQL Database Architecture")
    
    content_box = slide9.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Data Normative Tables & Security:"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(12)
    
    db_points = [
        ("Users Table", "Tracks user logins. Attributes: id (Primary Key), email (UNIQUE index), password_hash (SHA-256), phone, role (admin/user), and created_at timestamps."),
        ("Predictions Table", "Stores pathology history. Attributes: id (Primary Key), user_id (Foreign Key), crop, disease, confidence (Float), severity, image_path, and created_at."),
        ("DiseaseInformation Table", "Houses seeder recommendations. Attributes: crop, disease, description, symptoms, causes, organic_treatment, chemical_treatment, fertilizer, water_advice, prevention."),
        ("SQL Injection Prevention", "All inputs are strictly parameterized in Flask (e.g. cursor.execute('SELECT * FROM Users WHERE email = ?', (email,))), blocking malicious SQL injections.")
    ]
    
    for title, desc in db_points:
        p_db = tf.add_paragraph()
        p_db.text = f"✔ {title}: "
        p_db.font.name = "Arial"
        p_db.font.size = Pt(15)
        p_db.font.bold = True
        p_db.font.color.rgb = emerald_green
        
        p_desc = p_db.add_run()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = text_gray
        p_db.space_after = Pt(10)

    # Helper to insert screenshots into layout
    screenshots_dir = r"C:\Users\Administrator\Documents\crop-disease-diagnosis\backend\static\screenshots"
    
    # ------------------ SLIDE 10: SCREENSHOT - LOGIN ------------------
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide10)
    add_slide_header(slide10, "9. Interface Mockups: Access Gateway")
    
    left_box = slide10.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.0), Inches(5.0))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Authentication & Account Creation:"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(12)
    p2 = tf.add_paragraph()
    p2.text = "• Secure, glassmorphic card design.\n" \
              "• Includes inputs for username, phone, email, and secret passwords.\n" \
              "• Dynamically generates unique cartoon agricultural avatars on registration.\n" \
              "• Features real-time password strength indicators to enforce strong passwords."
    p2.font.name = "Arial"
    p2.font.size = Pt(15)
    p2.font.color.rgb = text_gray
    
    login_img = os.path.join(screenshots_dir, "login.jpg")
    if os.path.exists(login_img):
        # Image aspect ratio 16:9 -> fits nicely inside a box
        # login.jpg is ~1920x1080
        prs.slides[9].shapes.add_picture(login_img, Inches(6.0), Inches(1.8), Inches(6.5), Inches(4.5))

    # ------------------ SLIDE 11: SCREENSHOT - UPLOAD ------------------
    slide11 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide11)
    add_slide_header(slide11, "10. Interface Mockups: Diagnostic Laboratory")
    
    left_box = slide11.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.0), Inches(5.0))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Leaf Pathology Upload Portal:"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(12)
    p2 = tf.add_paragraph()
    p2.text = "• Interactive drag-and-drop workspace supporting leaf image uploads.\n" \
              "• Displays upload previews, loading states, and pipeline progress indicators.\n" \
              "• Guides the user visually with glowing cards and micro-animations."
    p2.font.name = "Arial"
    p2.font.size = Pt(15)
    p2.font.color.rgb = text_gray
    
    upload_img = os.path.join(screenshots_dir, "upload.jpg")
    if os.path.exists(upload_img):
        prs.slides[10].shapes.add_picture(upload_img, Inches(6.0), Inches(1.8), Inches(6.5), Inches(4.5))

    # ------------------ SLIDE 12: SCREENSHOT - PREDICTION ------------------
    slide12 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide12)
    add_slide_header(slide12, "11. Interface Mockups: Diagnosis Output")
    
    left_box = slide12.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.0), Inches(5.0))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Diagnostic Report & Pathology Summary:"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(12)
    p2 = tf.add_paragraph()
    p2.text = "• Displays diagnosis result, confidence scores, and crop health status.\n" \
              "• Details organic and chemical treatments, watering rules, and prevention guides.\n" \
              "• Includes inputs to customize Student Name and Project Title for the PDF report.\n" \
              "• Direct download triggers for university-grade PDF reports."
    p2.font.name = "Arial"
    p2.font.size = Pt(15)
    p2.font.color.rgb = text_gray
    
    pred_img = os.path.join(screenshots_dir, "prediction.jpg")
    if os.path.exists(pred_img):
        prs.slides[11].shapes.add_picture(pred_img, Inches(6.0), Inches(1.8), Inches(6.5), Inches(4.5))

    # ------------------ SLIDE 13: SCREENSHOT - DASHBOARD ------------------
    slide13 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide13)
    add_slide_header(slide13, "12. Interface Mockups: Admin Dashboard")
    
    left_box = slide13.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.0), Inches(5.0))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Control Metrics & Logs:"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(12)
    p2 = tf.add_paragraph()
    p2.text = "• Aggregates total predictions, accuracy levels, and crop distributions.\n" \
              "• Lists diagnostic history with date, submitter name, and confidence levels.\n" \
              "• Direct database sync controls and Excel spreadsheet download triggers."
    p2.font.name = "Arial"
    p2.font.size = Pt(15)
    p2.font.color.rgb = text_gray
    
    dash_img = os.path.join(screenshots_dir, "dashboard.jpg")
    if os.path.exists(dash_img):
        prs.slides[12].shapes.add_picture(dash_img, Inches(6.0), Inches(1.8), Inches(6.5), Inches(4.5))

    # ------------------ SLIDE 14: CONCLUSION ------------------
    slide14 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide14)
    add_slide_header(slide14, "13. Conclusion & Scope")
    
    content_box = slide14.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Key Accomplishments:"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(10)
    p2 = tf.add_paragraph()
    p2.text = "• Successfully built a full-stack crop disease diagnosis platform on port 8080.\n" \
              "• Integrated OpenCV color contours with neural networks to achieve consistent predictions.\n" \
              "• Compiled a self-executing Windows App (.exe) launcher for absolute offline portability."
    p2.font.name = "Arial"
    p2.font.size = Pt(15)
    p2.font.color.rgb = text_gray
    p2.space_after = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "Future Scope:"
    p3.font.name = "Arial"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = text_white
    p3.space_after = Pt(10)
    p4 = tf.add_paragraph()
    p4.text = "• Expand CNN training sets to support more agricultural pathogens (e.g. rice rusts, wheat smuts).\n" \
              "• Deploy model layers onto edge microcontrollers to run smart camera probes in irrigation fields.\n" \
              "• Integrate SMS notification gateways to alert farmers on local pathogen spreads automatically."
    p4.font.name = "Arial"
    p4.font.size = Pt(15)
    p4.font.color.rgb = text_gray

    # ------------------ SLIDE 15: THANK YOU & VIVA QUESTIONS ------------------
    slide15 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide15)
    
    # Title
    t_box = slide15.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1.2))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU!"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = emerald_green
    
    p2 = tf.add_paragraph()
    p2.text = "Review & Viva Questions Guidance"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = text_white
    p2.space_after = Pt(24)

    # Quick guide text
    guide_box = slide15.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9.333), Inches(2.5))
    tf = guide_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Viva Highlights to remember:"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "• Why Flask and SQLite? Flask acts as a fast gateway module, and SQLite provides a serverless relational database that runs locally on Windows without complex database configuration steps.\n" \
              "• Why OpenCV and MLP? OpenCV isolates true leaf pixels from shadow noise using HSV color masks, allowing a lightweight Multi-Layer Perceptron (MLP) network to calculate class scores rapidly on a CPU without requiring expensive CUDA GPUs."
    p2.font.name = "Arial"
    p2.font.size = Pt(14)
    p2.font.color.rgb = text_gray

    # Save presentation
    output_path = r"C:\Users\Administrator\Documents\crop-disease-diagnosis\Project_Presentation.pptx"
    prs.save(output_path)
    print("Project PowerPoint Presentation compiled successfully.")

if __name__ == "__main__":
    create_presentation()
